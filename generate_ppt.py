from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x06, 0x2B, 0x4C)
BLUE      = RGBColor(0x06, 0x5A, 0x82)
TEAL      = RGBColor(0x1C, 0x72, 0x93)
MINT      = RGBColor(0x02, 0xC3, 0x9A)
AMBER     = RGBColor(0xF5, 0xA6, 0x23)
CORAL     = RGBColor(0xF9, 0x61, 0x67)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xD6, 0xE8, 0xF5)
MGRAY     = RGBColor(0x64, 0x74, 0x8B)
DARK      = RGBColor(0x0F, 0x2A, 0x3F)
LBKG      = RGBColor(0xF0, 0xF7, 0xFF)
GREEN     = RGBColor(0x02, 0x8F, 0x76)

BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, line_color=None, line_w=0):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=13, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def bullets(slide, items, x, y, w, h, size=12, color=DARK, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            t, bld, col = item[0], item[1] if len(item)>1 else False, item[2] if len(item)>2 else color
        else:
            t, bld, col = item, (bold_first and i==0), color
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bld; r.font.color.rgb = col
    return tb

def footer(slide):
    rect(slide, 0, 7.1, 13.33, 0.4, BLUE)
    txt(slide, "IIT Mandi  |  Underwater Robotics and Autonomous Navigation",
        0.4, 7.12, 12.5, 0.3, size=10, color=LGRAY, align=PP_ALIGN.CENTER)

def banner(slide, phase_label, title, phase_color=BLUE):
    rect(slide, 0, 0, 13.33, 1.05, NAVY)
    txt(slide, phase_label, 0.4, 0.08, 8.0, 0.38,
        size=11, bold=True, color=phase_color)
    txt(slide, title, 0.4, 0.45, 12.5, 0.52, size=24, bold=True, color=WHITE)

def card(slide, x, y, w, h, fill=WHITE, border=LGRAY, bw=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bw)
    return s

def top_strip(slide, x, y, w, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.07))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.33, 7.5, NAVY)
rect(s1, 0, 0, 5.0, 7.5, BLUE)
rect(s1, 4.85, 0, 0.2, 7.5, TEAL)

txt(s1, "UNDERWATER ROBOTICS & AUTONOMOUS NAVIGATION",
    0.35, 0.4, 4.3, 0.6, size=10, bold=True, color=MINT)
txt(s1, "IIT Mandi", 0.35, 1.0, 4.0, 0.4, size=13, color=LGRAY)
lr = s1.shapes.add_shape(1, Inches(0.35), Inches(1.5), Inches(3.5), Inches(0.03))
lr.fill.solid(); lr.fill.fore_color.rgb = MINT; lr.line.fill.background()

# Phase badges
for i, (lbl, col) in enumerate([("PHASE 1", MINT), ("PHASE 2", AMBER)]):
    tag = s1.shapes.add_shape(1, Inches(0.35), Inches(1.75+i*0.62),
                               Inches(1.5), Inches(0.42))
    tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background()
    tf = tag.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = lbl
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY

txt(s1, "Dynamics Overview", 2.0, 1.78, 2.8, 0.38, size=11, color=WHITE)
txt(s1, "MATLAB Model (exOtter / Arjun)", 2.0, 2.4, 2.8, 0.38, size=11, color=WHITE)

txt(s1, "Presented by:", 0.35, 5.7, 4.0, 0.3, size=10, color=LGRAY)
txt(s1, "[Your Name]", 0.35, 6.0, 4.0, 0.45, size=15, bold=True, color=WHITE)
txt(s1, "[Roll No.]  |  [Branch]", 0.35, 6.5, 4.0, 0.35, size=10, color=LGRAY)

txt(s1, "Custom Control\nDeployment in\nArduSub", 5.3, 1.1, 7.7, 3.2,
    size=50, bold=True, color=WHITE)
txt(s1, "Assignment Progress Report — Phases 1 & 2",
    5.3, 4.5, 7.5, 0.5, size=15, italic=True, color=MINT)
txt(s1, "Based on study of exOtter.m (MSS Toolbox)  |  1-DOF Heave Simulation",
    5.3, 5.1, 7.5, 0.4, size=11, color=LGRAY)
footer(s1)

# =============================================================================
# SLIDE 2 — PHASE 1: What is Marine Craft Dynamics?
# (from your Phase 1 notes: Maneuvering vs Seakeeping, 6DOF intro)
# =============================================================================
s2 = prs.slides.add_slide(BLANK)
s2.background.fill.solid(); s2.background.fill.fore_color.rgb = LBKG
banner(s2, "PHASE 1  —  DYNAMICS OVERVIEW", "Marine Craft Motion — Two Ways Engineers Study It")

# Two big cards
for i, (title, col, content) in enumerate([
    ("1  Maneuvering", BLUE,
     ["How a ship steers and turns when the captain uses the engine/steering",
      "Studied by pretending the ship is in perfectly calm water",
      "Focuses on horizontal movements: forward, sliding sideways, turning",
      "Controlling the actual path the vessel takes"]),
    ("2  Seakeeping", TEAL,
     ["How a ship gets tossed around by ocean waves",
      "Assumes the ship is trying to maintain steady speed & straight course",
      "Studies vertical motions: heave, pitch, roll",
      "Understanding how waves affect the ship's stability"])
]):
    cx = 0.4 + i*6.5
    card(s2, cx, 1.2, 6.2, 5.55, border=col, bw=1.2)
    top_strip(s2, cx, 1.2, 6.2, col)
    hb = s2.shapes.add_shape(1, Inches(cx), Inches(1.27), Inches(6.2), Inches(0.6))
    hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
    txt(s2, title, cx+0.15, 1.3, 5.8, 0.52, size=15, bold=True, color=col)
    y = 2.0
    for pt in content:
        dot = s2.shapes.add_shape(9, Inches(cx+0.2), Inches(y+0.1), Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        txt(s2, pt, cx+0.42, y, 5.6, 0.55, size=12, color=DARK)
        y += 0.62

txt(s2, "These two categories together cover all six degrees of freedom of a vessel.",
    0.4, 6.88, 12.5, 0.3, size=11, italic=True, color=BLUE)
footer(s2)

# =============================================================================
# SLIDE 3 — PHASE 1: The Governing Equation (from your notes)
# M·v̇ + C(v)·v + D(v)·v + g(η) = τ
# =============================================================================
s3 = prs.slides.add_slide(BLANK)
s3.background.fill.solid(); s3.background.fill.fore_color.rgb = LBKG
banner(s3, "PHASE 1  —  DYNAMICS OVERVIEW", "The Key Equation  —  M·v̇ + C(v)·v + D(v)·v + g(η) = τ")

# Equation box
eb = s3.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.85))
eb.fill.solid(); eb.fill.fore_color.rgb = NAVY; eb.line.fill.background()
tf = eb.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "M · v̇   +   C(v) · v   +   D(v) · v   +   g(η)   =   τ"
r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = MINT; r.font.name = "Cambria"

# 5 term cards — from your notes
terms = [
    ("M", "Inertia Matrix", BLUE,
     ["M = Mrb + Ma", "Rigid body mass + Added mass", "From your notes: M(Mass Matrix)"]),
    ("C(v)·v", "Coriolis Forces", TEAL,
     ["Rotational coupling", "C = Crb + Ca", "Twisting momentum of the heavy vessel rotating in 3D"]),
    ("D(v)·v", "Drag / Damping", GREEN,
     ["Water resistance", "Non-linear at higher speeds", "D(v)·v — depends on both position and velocity"]),
    ("g(η)", "Gravity & Buoyancy", RGBColor(0x0C,0x44,0x7C),
     ["Natural desire to float", "Net restoring force", "Weight and distance from the water surface"]),
    ("τ", "Thruster Force", AMBER,
     ["Control input", "τ = [tau-X, tau-N]", "Force and torque from propellers"]),
]
cw = 2.36
for i, (sym, name, col, pts) in enumerate(terms):
    cx = 0.5 + i*(cw+0.09)
    card(s3, cx, 2.2, cw, 4.9, border=LGRAY)
    top_strip(s3, cx, 2.2, cw, col)
    hs = s3.shapes.add_shape(1, Inches(cx), Inches(2.27), Inches(cw), Inches(0.62))
    hs.fill.solid(); hs.fill.fore_color.rgb = col; hs.line.fill.background()
    txt(s3, sym, cx+0.05, 2.3, cw-0.1, 0.55,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s3, name, cx+0.1, 2.95, cw-0.2, 0.45,
        size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    y = 3.48
    for pt in pts:
        txt(s3, pt, cx+0.12, y, cw-0.24, 0.5, size=10.5, color=MGRAY)
        y += 0.5
footer(s3)

# =============================================================================
# SLIDE 4 — PHASE 1: Q1 Added Mass (from your notes)
# =============================================================================
s4 = prs.slides.add_slide(BLANK)
s4.background.fill.solid(); s4.background.fill.fore_color.rgb = LBKG
banner(s4, "PHASE 1  —  Q1: Added Mass", "What is Added Mass and Why Does It Matter More Underwater?")

# Definition card (left)
card(s4, 0.4, 1.2, 5.9, 5.55, border=BLUE, bw=1)
top_strip(s4, 0.4, 1.2, 5.9, BLUE)
txt(s4, "Definition", 0.6, 1.35, 5.5, 0.4, size=13, bold=True, color=BLUE)
txt(s4, "When a vehicle accelerates through water, it has to drag a surrounding blob of water with it. This extra water resists being moved — making the vehicle feel heavier than it actually is.\n\nFrom your notes: \"Added mass is the KE of the fluid that is forced to move along with the vehicle\"", 0.6, 1.82, 5.5, 2.5, size=12, color=DARK)

txt(s4, "From your notes:", 0.6, 4.35, 5.5, 0.35, size=11, bold=True, color=BLUE)
for pt in ["Physically, the fluid is added to the vehicle",
           "At low speeds: oscillates, difficult to move",
           "At high speed: becomes chaotic, multiplies aggressively"]:
    pass

for j, pt in enumerate(["Physically, the fluid is added to the vehicle",
           "At low speeds: oscillates, difficult to move",
           "At high speed: becomes chaotic, multiplies aggressively"]):
    dot = s4.shapes.add_shape(9, Inches(0.62), Inches(4.72+j*0.42), Inches(0.11), Inches(0.11))
    dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.fill.background()
    txt(s4, pt, 0.82, 4.68+j*0.42, 5.2, 0.4, size=11, color=DARK)

# Right card — comparison
card(s4, 6.55, 1.2, 6.38, 2.5, border=TEAL, bw=1)
top_strip(s4, 6.55, 1.2, 6.38, TEAL)
txt(s4, "Aircraft vs Underwater Vehicle", 6.75, 1.35, 6.0, 0.4, size=13, bold=True, color=TEAL)
table_data = [
    ("Air density", "≈ 1.2 kg/m³", "Added mass negligible"),
    ("Water density", "≈ 1025 kg/m³", "Added mass = 30–50% of vehicle mass"),
]
for j, (item, val, impact) in enumerate(table_data):
    ry = 1.82 + j*0.8
    txt(s4, item, 6.75, ry, 2.0, 0.38, size=11, bold=True, color=TEAL)
    txt(s4, val, 8.85, ry, 1.8, 0.38, size=11, color=DARK)
    txt(s4, impact, 10.72, ry, 2.05, 0.38, size=10.5, italic=True, color=MGRAY)

# Why it matters card
card(s4, 6.55, 3.9, 6.38, 2.85, border=AMBER, bw=1)
top_strip(s4, 6.55, 3.9, 6.38, AMBER)
txt(s4, "Why It Matters for AUVs", 6.75, 4.05, 6.0, 0.4, size=13, bold=True, color=AMBER)
for j, pt in enumerate([
    "Makes vehicle feel much heavier than it is",
    "Added mass must be included in M matrix: M = Mrb + Ma",
    "Bigger mass → harder to move the boat",
    "Smaller mass → easier to accelerate",
    "From notes: A33 = 5.0 kg for heave direction"
]):
    dot = s4.shapes.add_shape(9, Inches(6.75), Inches(4.55+j*0.42), Inches(0.11), Inches(0.11))
    dot.fill.solid(); dot.fill.fore_color.rgb = AMBER; dot.line.fill.background()
    txt(s4, pt, 6.95, 4.51+j*0.42, 5.8, 0.4, size=11, color=DARK)

footer(s4)

# =============================================================================
# SLIDE 5 — PHASE 1: Q2 Nonlinear Damping + Q3 6-DOF
# =============================================================================
s5 = prs.slides.add_slide(BLANK)
s5.background.fill.solid(); s5.background.fill.fore_color.rgb = LBKG
banner(s5, "PHASE 1  —  Q2 & Q3", "Nonlinear Damping & The 6 Degrees of Freedom")

# Q2 left
card(s5, 0.4, 1.2, 6.0, 5.55, border=TEAL, bw=1)
top_strip(s5, 0.4, 1.2, 6.0, TEAL)
txt(s5, "Q2  Why is Damping Nonlinear at Higher Speeds?", 0.6, 1.32, 5.7, 0.55,
    size=13, bold=True, color=TEAL)
txt(s5, "Linear drag (low speed):", 0.6, 1.98, 5.5, 0.35, size=11, bold=True, color=TEAL)
txt(s5, "F_drag = D · v      (proportional to velocity)", 0.7, 2.35, 5.3, 0.4, size=12, color=DARK)
txt(s5, "Nonlinear drag (high speed):", 0.6, 2.85, 5.5, 0.35, size=11, bold=True, color=CORAL)
txt(s5, "F_drag = D · v²     (proportional to velocity squared)", 0.7, 3.22, 5.3, 0.4, size=12, color=DARK)

txt(s5, "Your notes explanation:", 0.6, 3.75, 5.5, 0.35, size=11, bold=True, color=TEAL)
for j, pt in enumerate([
    "Damping becomes non-linear because at higher speeds, the fluid flow transitions from smooth (laminar) to turbulent",
    "Drag grows much faster — the vehicle loses energy to the fluid rapidly",
    "At high speed: the fluid is highly dynamic and produces chaotic drag",
    "Linear model is only a reasonable approximation at small/slow motions"
]):
    txt(s5, pt, 0.65, 4.18+j*0.52, 5.6, 0.5, size=11, color=DARK)

# Q3 right
card(s5, 6.6, 1.2, 6.33, 5.55, border=BLUE, bw=1)
top_strip(s5, 6.6, 1.2, 6.33, BLUE)
txt(s5, "Q3  The 6 Degrees of Freedom", 6.8, 1.32, 6.0, 0.45, size=13, bold=True, color=BLUE)

dofs = [
    ("Surge", "Forward / Backward (X)", BLUE, False),
    ("Sway",  "Left / Right (Y)", BLUE, False),
    ("Heave", "Up / Down — DEPTH (Z)", MINT, True),
    ("Roll",  "Rotation about X", TEAL, False),
    ("Pitch", "Rotation about Y", TEAL, False),
    ("Yaw",   "Rotation about Z — HEADING", AMBER, True),
]
for j, (name, desc, col, star) in enumerate(dofs):
    ry = 1.88 + j*0.6
    rb = s5.shapes.add_shape(1, Inches(6.75), Inches(ry), Inches(6.0), Inches(0.52))
    rb.fill.solid()
    rb.fill.fore_color.rgb = NAVY if star else WHITE
    rb.line.color.rgb = col; rb.line.width = Pt(1.2 if star else 0.5)
    tab = s5.shapes.add_shape(1, Inches(6.75), Inches(ry), Inches(0.07), Inches(0.52))
    tab.fill.solid(); tab.fill.fore_color.rgb = col; tab.line.fill.background()
    tcol = MINT if star else col
    txt(s5, name, 6.92, ry+0.07, 1.2, 0.38, size=12, bold=True, color=tcol)
    txt(s5, desc, 8.18, ry+0.07, 3.5, 0.38, size=11, color=WHITE if star else DARK)
    if star:
        txt(s5, "KEY", 11.9, ry+0.07, 0.8, 0.38, size=10, bold=True, color=AMBER)

txt(s5, "Heave (depth) and Yaw (heading) are the two most important DOF for AUV control",
    6.8, 6.56, 6.1, 0.38, size=11, italic=True, color=AMBER)
footer(s5)

# =============================================================================
# SLIDE 6 — PHASE 2: exOtter / What I Studied
# =============================================================================
s6 = prs.slides.add_slide(BLANK)
s6.background.fill.solid(); s6.background.fill.fore_color.rgb = LBKG
banner(s6, "PHASE 2  —  MATLAB MODEL", "What I Studied: exOtter.m from MSS Toolbox", AMBER)

# Left: what exOtter is
card(s6, 0.4, 1.2, 6.1, 5.55, border=AMBER, bw=1)
top_strip(s6, 0.4, 1.2, 6.1, AMBER)
txt(s6, "exOtter.m — The Reference Model I Studied", 0.6, 1.32, 5.8, 0.45,
    size=13, bold=True, color=AMBER)
txt(s6, "From my notes (Phase 2, Page 1):", 0.6, 1.88, 5.7, 0.35,
    size=11, bold=True, color=AMBER)
for j, pt in enumerate([
    "Ex-1 = exOtter.m  →  Simulates a boat (Otter USV) using MSS Toolbox",
    "Input: Propeller forces (tau)",
    "Output: Motion (velocity, position)",
    "Simulation runs at 50 Hz  →  h = 1/50 = 0.02 sec time step",
    "Uses EKF (Extended Kalman Filter) for state estimation",
    "States: u, v, w (translational velocities), p, q, r (angular velocities)",
    "phi, theta → angular positions  |  psi → heading",
]):
    dot = s6.shapes.add_shape(9, Inches(0.62), Inches(2.33+j*0.52), Inches(0.11), Inches(0.11))
    dot.fill.solid(); dot.fill.fore_color.rgb = AMBER; dot.line.fill.background()
    txt(s6, pt, 0.82, 2.3+j*0.52, 5.5, 0.5, size=11, color=DARK)

# Right: key variables from code
card(s6, 6.7, 1.2, 6.33, 2.55, border=BLUE, bw=1)
top_strip(s6, 6.7, 1.2, 6.33, BLUE)
txt(s6, "Key Variables in exOtter.m", 6.9, 1.32, 6.0, 0.42, size=13, bold=True, color=BLUE)
code_lines = [
    "x = [u v w p q r x y z phi theta psi]'",
    "B_prop  →  Input (propeller) matrix",
    "Binv = inv(B_prop)  →  Control allocation",
    "M(6,6)  →  Mass matrix element",
    "EKF_5states  →  Estimates [x_N, v_N, U, chi, omega]",
]
for j, line in enumerate(code_lines):
    txt(s6, line, 6.9, 1.85+j*0.42, 5.9, 0.4,
        size=11, color=DARK if j>0 else BLUE)

# Binv explanation
card(s6, 6.7, 3.9, 6.33, 2.85, border=TEAL, bw=1)
top_strip(s6, 6.7, 3.9, 6.33, TEAL)
txt(s6, "What Binv Means (from your notes)", 6.9, 4.02, 6.0, 0.42, size=13, bold=True, color=TEAL)
txt(s6, "Binv = inv(B_prop)  →  action(C)", 6.9, 4.52, 6.0, 0.38, size=12, color=DARK)
for j, pt in enumerate([
    "This is the mathematical model",
    "We use this to calculate required force",
    "Instead of requiring forces → base → required propeller speed",
    "Binv links the problem: force needed → motor command"
]):
    txt(s6, pt, 6.9, 4.98+j*0.42, 6.0, 0.4, size=11, color=DARK)

footer(s6)

# =============================================================================
# SLIDE 7 — PHASE 2: PID Controller 
# =============================================================================
s7 = prs.slides.add_slide(BLANK)
s7.background.fill.solid(); s7.background.fill.fore_color.rgb = LBKG
banner(s7, "PHASE 2  —  MATLAB MODEL", "PID Controller — Pole Placement Method", AMBER)

# Equation banner
eq = s7.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.75))
eq.fill.solid(); eq.fill.fore_color.rgb = NAVY; eq.line.fill.background()
tf = eq.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "tau_N = (T/K)·a_d + (1/K)·omega_d  −  Kp·[ ssa(chi_hat − chi_d) + Td·(omega_hat − omega_d) + (1/Ti)·z ]"
r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = MINT; r.font.name = "Consolas"

# 3 PID term cards
pid_terms = [
    ("P — Proportional", BLUE,
     "Kp = M(6,6) × wn²",
     ["Controls how strongly the system reacts to error (chi_hat - chi_d)",
      "Bigger Kp → reacts faster but may overshoot",
      "Your notes: Kp controls how hard to push",
      "Big error → big correction immediately"]),
    ("D — Derivative", TEAL,
     "Kd = M(6,6) × (2×zeta×wn − 1/T)\nTd = Kd / Kp",
     ["Adds damping — slows down if approaching too fast",
      "Looks at rate of change of error (omega_hat − omega_d)",
      "Your notes: Kd adds slowing",
      "Prevents overshoot and oscillation"]),
    ("I — Integral", GREEN,
     "Ti = 10 / wn\nz = z + h × ssa(chi_hat − chi_d)",
     ["Fixes small long-term errors that keep happening",
      "From your notes: 'Integral fixes steady-state error'",
      "Small errors accumulate → z builds up → fixes it",
      "z = 0 at start, then it stores accumulated errors"]),
]
cw = 4.0
for i, (title, col, formula, pts) in enumerate(pid_terms):
    cx = 0.4 + i*(cw+0.08)
    card(s7, cx, 2.1, cw, 4.65, border=col, bw=1)
    top_strip(s7, cx, 2.1, cw, col)
    hb = s7.shapes.add_shape(1, Inches(cx), Inches(2.17), Inches(cw), Inches(0.5))
    hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
    txt(s7, title, cx+0.12, 2.2, cw-0.2, 0.44, size=13, bold=True, color=col)
    # formula box
    fb = s7.shapes.add_shape(1, Inches(cx+0.12), Inches(2.75), Inches(cw-0.24), Inches(0.6))
    fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0xE0,0xF0,0xFF)
    fb.line.color.rgb = col; fb.line.width = Pt(0.5)
    txt(s7, formula, cx+0.2, 2.78, cw-0.4, 0.55, size=10.5, color=DARK)
    for j, pt in enumerate(pts):
        dot = s7.shapes.add_shape(9, Inches(cx+0.18), Inches(3.45+j*0.5), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        txt(s7, pt, cx+0.36, 3.42+j*0.5, cw-0.5, 0.48, size=11, color=DARK)

footer(s7)

# =============================================================================
# SLIDE 8 — PHASE 2: Reference Model + Tuning
# =============================================================================
s8 = prs.slides.add_slide(BLANK)
s8.background.fill.solid(); s8.background.fill.fore_color.rgb = LBKG
banner(s8, "PHASE 2  —  MATLAB MODEL", "Reference Model, Tuning & Propeller Dynamics", AMBER)

# Left: Reference model (wn, zeta, your notes)
card(s8, 0.4, 1.2, 6.0, 5.55, border=BLUE, bw=1)
top_strip(s8, 0.4, 1.2, 6.0, BLUE)
txt(s8, "Reference Model — Why We Need It", 0.6, 1.32, 5.7, 0.45, size=13, bold=True, color=BLUE)
txt(s8, "Without a reference model → Sudden jump → difficult to control\nReference model is like a filter — smooth input → smooth desired motion",
    0.6, 1.85, 5.6, 0.72, size=11, color=DARK)

txt(s8, "Parameters (from your notes):", 0.6, 2.65, 5.5, 0.35, size=11, bold=True, color=BLUE)
params = [
    ("wn_d = 1.0", "Natural frequency of reference system"),
    ("zeta_d = 1.0", "Damping — Controls oscillation smoothness"),
    ("r_max = deg2rad(10)", "Maximum turning rate limit"),
    ("chi_d = 0", "Desired course direction (starts at 0)"),
    ("omega_d = 0", "Desired course rate (starts at 0)"),
    ("a_d = 0", "Desired course acceleration (starts at 0)"),
]
for j, (param, desc) in enumerate(params):
    txt(s8, param, 0.65, 3.08+j*0.42, 2.2, 0.38, size=11, bold=True, color=BLUE)
    txt(s8, desc,  2.92, 3.08+j*0.42, 3.3, 0.38, size=11, color=DARK)

txt(s8, "Zeta (damping):", 0.6, 5.7, 5.5, 0.3, size=11, bold=True, color=TEAL)
txt(s8, "zeta < 1 → oscillates  |  zeta = 1 → perfectly smooth  |  zeta > 1 → very close, no oscillation",
    0.6, 6.05, 5.6, 0.38, size=10.5, italic=True, color=DARK)

# Right: Propeller dynamics + load condition
card(s8, 6.6, 1.2, 6.33, 2.6, border=TEAL, bw=1)
top_strip(s8, 6.6, 1.2, 6.33, TEAL)
txt(s8, "Propeller Dynamics (from your notes)", 6.8, 1.32, 6.0, 0.42, size=13, bold=True, color=TEAL)
for j, (param, desc) in enumerate([
    ("Tn = 0.1 s", "Propeller time constant — how fast propeller responds"),
    ("n = [0 0]'", "Initial propeller speeds: left=0, right=0"),
    ("n_c = sign(u)×sqrt(|u|)", "Desired propeller speed from control"),
    ("n updated by Euler:", "n = n − h/Tn × (n − n_c)"),
]):
    txt(s8, param, 6.8, 1.85+j*0.42, 2.2, 0.38, size=11, bold=True, color=TEAL)
    txt(s8, desc, 9.1, 1.85+j*0.42, 3.7, 0.38, size=11, color=DARK)

card(s8, 6.6, 3.95, 6.33, 2.8, border=AMBER, bw=1)
top_strip(s8, 6.6, 3.95, 6.33, AMBER)
txt(s8, "Load Condition & Ocean Current", 6.8, 4.07, 6.0, 0.42, size=13, bold=True, color=AMBER)
for j, (param, desc) in enumerate([
    ("mp = 25 kg", "Payload mass (max 45 kg) — extra weight on boat"),
    ("rp = [0.05 0 -0.35]'", "Location of payload on the boat"),
    ("V_c = 0", "Current speed — set to 0 (no current)"),
    ("beta_c = 30° (π/180)", "Current direction — 30 degrees"),
    ("Understanding:", "This action defines extra weight on the boat and the environment that can affect its motion"),
]):
    txt(s8, param, 6.8, 4.6+j*0.42, 2.4, 0.38, size=11, bold=True, color=AMBER)
    txt(s8, desc, 9.28, 4.6+j*0.42, 3.5, 0.38, size=11, color=DARK)

footer(s8)

# =============================================================================
# SLIDE 9 — PHASE 2: EKF + simdata 
# =============================================================================
s9 = prs.slides.add_slide(BLANK)
s9.background.fill.solid(); s9.background.fill.fore_color.rgb = LBKG
banner(s9, "PHASE 2  —  MATLAB MODEL", "State Estimation: EKF + simdata Table", AMBER)

# EKF left
card(s9, 0.4, 1.2, 6.1, 5.55, border=TEAL, bw=1)
top_strip(s9, 0.4, 1.2, 6.1, TEAL)
txt(s9, "EKF — Extended Kalman Filter (State Estimation)", 0.6, 1.32, 5.8, 0.45,
    size=13, bold=True, color=TEAL)
txt(s9, "From your notes:", 0.6, 1.85, 5.5, 0.32, size=11, bold=True, color=TEAL)
for j, pt in enumerate([
    "EKF estimates the state — it does NOT control the boat",
    "Uses GPS sensor data to estimate where the boat is",
    "x_hat = [x_N, v_N, U, chi, omega]'  →  5 states",
    "1→ estimated x position (North)",
    "2→ estimated y position (v_N)",
    "3→ estimated speed (U)",
    "4→ estimated course angle (chi)",
    "5→ estimated turning rate (omega)",
]):
    dot = s9.shapes.add_shape(9, Inches(0.62), Inches(2.28+j*0.48), Inches(0.11), Inches(0.11))
    dot.fill.solid(); dot.fill.fore_color.rgb = TEAL; dot.line.fill.background()
    txt(s9, pt, 0.82, 2.25+j*0.48, 5.5, 0.46, size=11, color=DARK)

txt(s9, "Covariance matrices from your notes:", 0.6, 6.45, 5.5, 0.35,
    size=10.5, bold=True, color=TEAL)
txt(s9, "Qd = 500 × diag([1000 1000])  |  Rd = 0.00000001 × diag([1,1])",
    0.6, 6.82, 5.5, 0.35, size=10, italic=True, color=MGRAY)

# Right: simdata explanation
card(s9, 6.65, 1.2, 6.38, 2.7, border=BLUE, bw=1)
top_strip(s9, 6.65, 1.2, 6.38, BLUE)
txt(s9, "simdata Table — How Data is Stored", 6.85, 1.32, 6.0, 0.42, size=13, bold=True, color=BLUE)
txt(s9, "simdata = zeros(length(t), 19)", 6.85, 1.83, 6.0, 0.38, size=11, bold=True, color=DARK)
txt(s9, "From your notes: 'Big purposeful data saves all important values into a table'\nRow 1 = first time step data  |  Row 2 = time step 2 data",
    6.85, 2.25, 6.1, 0.65, size=11, color=DARK)
for j, (col_range, what) in enumerate([
    ("Cols 1–12:", "x = [u,v,w,p,q,r,x,y,z,phi,theta,psi]"),
    ("Cols 13–17:", "x_hat from EKF (5 states)"),
    ("Cols 18–19:", "n = [n1, n2] propeller speeds"),
]):
    txt(s9, col_range, 6.85, 3.0+j*0.42, 1.7, 0.38, size=11, bold=True, color=BLUE)
    txt(s9, what, 8.62, 3.0+j*0.42, 4.3, 0.38, size=11, color=DARK)

# Right bottom: RK4
card(s9, 6.65, 4.05, 6.38, 2.7, border=GREEN, bw=1)
top_strip(s9, 6.65, 4.05, 6.38, GREEN)
txt(s9, "RK4 — Runge-Kutta 4th Order Integration", 6.85, 4.17, 6.0, 0.42,
    size=13, bold=True, color=GREEN)
txt(s9, "x = rk4(@otter, h, x, n, mp, rp, V_c, beta_c)", 6.85, 4.68, 6.0, 0.38,
    size=10.5, bold=True, color=DARK)
for j, pt in enumerate([
    "RK4 is a numerical method to solve differential equations",
    "Better accuracy than Euler (4 slope estimates per step)",
    "From your notes: 'It is almost mini-calculations using the estimate motion using 4 directions'",
    "Used for main vehicle dynamics — Euler used for simpler updates (z, n)"
]):
    txt(s9, pt, 6.85, 5.15+j*0.38, 6.0, 0.36, size=10.5, color=DARK)

footer(s9)

# =============================================================================
# SLIDE 10 — PHASE 2: My 1-DOF Model (Arjun)
# =============================================================================
s10 = prs.slides.add_slide(BLANK)
s10.background.fill.solid(); s10.background.fill.fore_color.rgb = LBKG
banner(s10, "PHASE 2  —  MY OWN MODEL", "heave_1DOF.m  —  My 1-DOF Heave Simulation", AMBER)

# Equation
eq2 = s10.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.75))
eq2.fill.solid(); eq2.fill.fore_color.rgb = NAVY; eq2.line.fill.background()
tf2 = eq2.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
r2 = tf2.paragraphs[0].add_run()
r2.text = "(m + A₃₃) · ẇ   +   D₃₃ · w   =   τ_heave"
r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = MINT; r2.font.name = "Cambria"

# My parameters
card(s10, 0.4, 2.1, 5.5, 5.15, border=AMBER, bw=1)
top_strip(s10, 0.4, 2.1, 5.5, AMBER)
txt(s10, "My Model Parameters", 0.6, 2.22, 5.2, 0.42, size=13, bold=True, color=AMBER)
my_params = [
    ("m = 11.5 kg", "Vehicle mass"),
    ("A33 = 5.0 kg", "Added mass in heave direction"),
    ("D33 = 25.0 kg/s", "Linear damping coefficient"),
    ("M_heave = m + A33 = 16.5 kg", "Total effective inertia"),
    ("h = 0.01 s", "Time step (Euler integration)"),
    ("T_final = 20 s", "Simulation duration"),
    ("t_step = 2 s", "Force switches ON at 2 seconds"),
    ("F_step = 10 N", "Step thruster force applied"),
]
for j, (param, desc) in enumerate(my_params):
    txt(s10, param, 0.6, 2.73+j*0.42, 2.5, 0.38, size=11, bold=True, color=AMBER)
    txt(s10, desc, 3.18, 2.73+j*0.42, 2.6, 0.38, size=11, color=DARK)

# The loop
card(s10, 6.1, 2.1, 6.83, 5.15, border=TEAL, bw=1)
top_strip(s10, 6.1, 2.1, 6.83, TEAL)
txt(s10, "The Euler Loop — Heart of the Simulation", 6.3, 2.22, 6.5, 0.42,
    size=13, bold=True, color=TEAL)
loop_steps = [
    ("Step 1 — Acceleration:", "w_dot = (tau_heave − D33×w) / M_heave\nNewton's 2nd law: a = F_net / mass"),
    ("Step 2 — Rate of depth:", "z_dot = w(i)\nDepth changes at same rate as velocity"),
    ("Step 3 — Euler update:", "w(i+1) = w(i) + h × w_dot\nz(i+1) = z(i) + h × z_dot"),
    ("Step 4 — Analytical check:", "w_ss = tau / D33 = 10/25 = 0.4 m/s\nTerminal velocity — confirms simulation correct"),
]
for j, (title, code) in enumerate(loop_steps):
    ry = 2.73 + j*1.1
    lb = s10.shapes.add_shape(1, Inches(6.25), Inches(ry), Inches(6.5), Inches(1.0))
    lb.fill.solid(); lb.fill.fore_color.rgb = LBKG; lb.line.color.rgb = TEAL; lb.line.width = Pt(0.5)
    txt(s10, title, 6.4, ry+0.05, 6.3, 0.35, size=11, bold=True, color=TEAL)
    txt(s10, code, 6.4, ry+0.42, 6.3, 0.55, size=10.5, color=DARK)

footer(s10)

# =============================================================================
# SLIDE 11 — PHASE 2: Results (3 plots)
# =============================================================================
s11 = prs.slides.add_slide(BLANK)
s11.background.fill.solid(); s11.background.fill.fore_color.rgb = LBKG
banner(s11, "PHASE 2  —  RESULTS", "Simulation Output — 3 Plots from heave_1DOF.m", AMBER)

plots = [
    ("Plot 1\nDepth vs Time", BLUE,
     "[INSERT MATLAB PLOT 1 HERE\n— depth (m) vs time (s)]",
     ["Depth = 0 m for first 2 s",
      "After t=2s: depth increases continuously",
      "No controller → depth grows without bound",
      "This motivates Phase 3: PID control"]),
    ("Plot 2\nHeave Velocity vs Time", TEAL,
     "[INSERT MATLAB PLOT 2 HERE\n— velocity (m/s) vs time (s)]",
     ["Starts at 0 m/s at rest",
      "Rises rapidly after force applied at t=2s",
      "Levels off at terminal: w = 0.4 m/s",
      "w_ss = tau/D33 = 10/25 = 0.4 m/s"]),
    ("Plot 3\nThruster Force vs Time", AMBER,
     "[INSERT MATLAB PLOT 3 HERE\n— force (N) vs time (s)]",
     ["Zero force for t < 2 s",
      "Step jump to 10 N at t = 2 s",
      "Constant — this is open-loop",
      "Confirms step input applied correctly"]),
]
pw = 4.1
for i, (title, col, placeholder, pts) in enumerate(plots):
    cx = 0.35 + i*(pw+0.1)
    card(s11, cx, 1.2, pw, 5.65, border=col, bw=1)
    hb = s11.shapes.add_shape(1, Inches(cx), Inches(1.2), Inches(pw), Inches(0.55))
    hb.fill.solid(); hb.fill.fore_color.rgb = col; hb.line.fill.background()
    txt(s11, title, cx+0.12, 1.22, pw-0.2, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ip = s11.shapes.add_shape(1, Inches(cx+0.12), Inches(1.83), Inches(pw-0.24), Inches(2.6))
    ip.fill.solid(); ip.fill.fore_color.rgb = RGBColor(0xE4,0xF2,0xFF)
    ip.line.color.rgb = col; ip.line.width = Pt(0.5)
    txt(s11, placeholder, cx+0.22, 2.65, pw-0.44, 0.8,
        size=10.5, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)
    txt(s11, "Screenshot from MATLAB here", cx+0.22, 2.0, pw-0.44, 0.4,
        size=10, color=col, align=PP_ALIGN.CENTER)
    for j, pt in enumerate(pts):
        dot = s11.shapes.add_shape(9, Inches(cx+0.18), Inches(4.54+j*0.38), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        txt(s11, pt, cx+0.36, 4.5+j*0.38, pw-0.5, 0.36, size=10.5, color=DARK)

footer(s11)

# =============================================================================
# SLIDE 12 — SUMMARY & NEXT STEPS
# =============================================================================
s12 = prs.slides.add_slide(BLANK)
rect(s12, 0, 0, 13.33, 7.5, NAVY)
rect(s12, 0, 0, 0.22, 7.5, MINT)

txt(s12, "Summary & What I Learned", 0.5, 0.3, 10.0, 0.75,
    size=32, bold=True, color=WHITE)
txt(s12, "Phases 1 & 2 — Based on my own study notes",
    0.5, 1.05, 8.0, 0.42, size=14, italic=True, color=MINT)

for i, (phase, col, items) in enumerate([
    ("PHASE 1 — Dynamics Overview", MINT, [
        "Studied two ways engineers analyze ships: Maneuvering vs Seakeeping",
        "Understood the full 6-DOF equation: M·v̇ + C(v)·v + D(v)·v + g(η) = τ",
        "Added Mass: vehicle drags water → feels heavier (A33 = 30-50% of mass)",
        "Nonlinear damping: turbulent flow at high speed → drag grows as v²",
        "6 DOF identified: Heave (depth) and Yaw (heading) are most critical",
        "Understood M matrix, Coriolis, Drag, Gravity/Buoyancy, and τ",
    ]),
    ("PHASE 2 — MATLAB Model", AMBER, [
        "Studied exOtter.m from MSS Toolbox in detail (EKF, PID, reference model)",
        "Understood EKF: estimates [x_N, v_N, U, chi, omega] from GPS data",
        "Understood PID pole placement: Kp, Kd, Ti, reference model wn_d=1, zeta_d=1",
        "Built my own heave_1DOF.m — 1-DOF depth simulation from scratch",
        "Applied 10N step force: vehicle sinks, velocity → terminal 0.4 m/s",
        "Verified: w_ss = tau/D33 = 10/25 = 0.4 m/s (analytical = simulated)",
    ]),
]):
    cb = s12.shapes.add_shape(1, Inches(0.4+i*6.45), Inches(1.6),
                               Inches(6.25), Inches(4.25))
    cb.fill.solid(); cb.fill.fore_color.rgb = RGBColor(0x07,0x35,0x60)
    cb.line.fill.background()
    top_strip(s12, 0.4+i*6.45, 1.6, 6.25, col)
    txt(s12, phase, 0.6+i*6.45, 1.72, 6.0, 0.42, size=13, bold=True, color=col)
    for j, item in enumerate(items):
        txt(s12, "✓  " + item, 0.6+i*6.45, 2.22+j*0.52, 5.95, 0.5, size=11, color=WHITE)

# Next steps
nb = s12.shapes.add_shape(1, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.88))
nb.fill.solid(); nb.fill.fore_color.rgb = RGBColor(0x02,0x39,0x5A); nb.line.fill.background()
txt(s12, "Coming Next:", 0.6, 6.1, 2.0, 0.38, size=12, bold=True, color=MINT)
for k, nx in enumerate(["Phase 3: PID\nDepth Controller",
                          "Phase 4: ArduSub\nCodebase Map",
                          "Phase 5: Custom\nUserCode.cpp",
                          "Phase 6: SITL\nBuild & Test"]):
    ax = 2.6 + k*2.5
    ab = s12.shapes.add_shape(1, Inches(ax), Inches(6.08), Inches(2.3), Inches(0.72))
    ab.fill.solid(); ab.fill.fore_color.rgb = BLUE; ab.line.fill.background()
    txt(s12, nx, ax+0.08, 6.1, 2.15, 0.68, size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

txt(s12, "Thank You", 0.5, 6.98, 12.33, 0.42,
    size=14, color=LGRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\harsh\.antigravity\uav-portfolio\ArduSub_FromNotes.pptx"
prs.save(out)
print("Saved:", out)
